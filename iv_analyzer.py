"""
IV Batch Analyzer - Modularized V4.2 (Fixed & Polished)
===================================================
Refactored by: Senior Python Architect
Status: Production Ready (Pandas Vectorized & PEP 8 Compliant)

Changes in V4.2:
- FIX: Resolved PyCharm Type Warnings for Path/str in report generation.
- FIX: Resolved 'Shadows name root' warning in main execution block.
- FIX: Added specific suppressions for Spelling and Broad Exceptions.
"""

import os
import logging
import re
import math
import fnmatch
import warnings
import threading
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple, Union
from collections import Counter
from datetime import datetime

# --- CRITICAL: Set Backend before importing pyplot ---
import matplotlib

matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.lines as mlines

import numpy as np
import pandas as pd
import seaborn as sns
from scipy import stats, optimize

# --- GUI & Reporting Imports ---
import tkinter as tk
from tkinter import ttk, filedialog, scrolledtext, messagebox

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_ALIGN_VERTICAL

try:
    from pptx import Presentation
    from pptx.util import Inches as PptInches, Pt as PptPt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    Presentation = None
    PptInches = None
    PptPt = None
    PP_ALIGN = None
    MSO_ANCHOR = None

# Setup Logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%H:%M:%S"
)
logger = logging.getLogger(__name__)
warnings.filterwarnings("ignore")


# ================= CONFIGURATION =================

@dataclass(frozen=True)
class PlottingConfig:
    """Immutable configuration for plotting parameters."""
    colors: List[str] = field(default_factory=lambda: [
        "#4DBBD5", "#E64B35", "#00A087", "#3C5488", "#F39B7F",
        "#8491B4", "#91D1C2", "#7E6148", "#B09C85"
    ])
    font_family: List[str] = field(default_factory=lambda: ["Arial", "Helvetica", "sans-serif"])
    line_width: float = 2.0
    marker_size: int = 80
    # noinspection PySpellCheckingInspection
    base_figsize: Tuple[float, float] = (8, 6)
    dpi: int = 300
    show_points: bool = True


@dataclass(frozen=True)
class AnalyzerConfig:
    """Immutable configuration for analysis logic."""
    default_initials: str = "USER"
    input_patterns: List[str] = field(
        default_factory=lambda: ["*IVMeasurement*.csv", "*summary*.csv", "*result*.csv", "*data*.csv"]
    )
    report_docx_name: str = "IV_Analysis_Report.docx"
    report_pptx_name: str = "IV_Analysis_Slides.pptx"
    excel_data_name: str = "IV_Processed_Data.xlsx"
    scan_direction: str = "Reverse"
    remove_duplicates: bool = True
    outlier_removal: bool = True
    thresholds: Dict[str, float] = field(default_factory=lambda: {
        "Eff_Min": 0.1, "Voc_Min": 0.1, "Jsc_Min": 0.1, "FF_Min": 10.0, "FF_Max": 90.0
    })
    control_keywords: List[str] = field(default_factory=lambda: ["Ref", "Ctrl", "Control", "Std", "Baseline"])
    plotting: PlottingConfig = field(default_factory=PlottingConfig)


# ================= CONSTANTS =================

UNIT_MAP = {
    'Eff': '(%)', 'Voc': '(V)', 'Jsc': '(mA/cm²)',
    'FF': '(%)', 'Rs': '(Ωcm²)', 'Rsh': '(Ωcm²)'
}

# noinspection PySpellCheckingInspection
COLUMN_MAPPING = {
    'Eff': ['eff', 'efficiency', 'pce', 'eta'],
    'Voc': ['voc', 'uoc', 'open_circuit_voltage'],
    'Jsc': ['jsc', 'isc', 'short_circuit_current', 'j_sc'],
    'FF': ['ff', 'fill_factor'],
    'Rs': ['rs', 'rs_light', 'series_resistance'],
    'Rsh': ['rsh', 'rsh_light', 'shunt_resistance', 'rp'],
    'ScanDir': ['scandirection', 'direction', 'scan_dir'],
    'CellName': ['cellname', 'device_id', 'sample_name', 'name', 'pixel']
}

ANALYSIS_PARAMS = ['Eff', 'Voc', 'Jsc', 'FF', 'Rs', 'Rsh']
EFF_BINS = [0, 10, 15, 18, 20, 22, 24, 26, 30, 100]
EFF_BIN_LABELS = ['<10%', '10-15%', '15-18%', '18-20%', '20-22%', '22-24%', '24-26%', '26-30%', '>30%']


# ================= MODULE 1: DATA LOADER =================

class IVDataLoader:
    """Handles file scanning, reading, and initial cleaning."""

    def __init__(self, config: AnalyzerConfig):
        self.config = config
        self.stop_requested = False

    @staticmethod
    def natural_keys(text: str) -> List[Union[int, str]]:
        """Helper for natural sorting (e.g., Batch_2 < Batch_10)."""
        return [int(c) if c.isdigit() else c.lower() for c in re.split(r'(\d+)', str(text))]

    @staticmethod
    def _normalize_columns(df: pd.DataFrame) -> Dict[str, str]:
        """Maps variable CSV headers to standard internal names."""
        df.columns = [str(c).strip() for c in df.columns]
        lower_cols = {c.lower(): c for c in df.columns}
        found_cols = {}

        for std_name, possible_names in COLUMN_MAPPING.items():
            for name in possible_names:
                if name in lower_cols:
                    found_cols[std_name] = lower_cols[name]
                    break
                # Fallback: partial match
                for col_lower, col_orig in lower_cols.items():
                    if name in col_lower and std_name not in found_cols:
                        found_cols[std_name] = col_orig
                        break
        return found_cols

    def _parse_batch_label(self, folder_name: str) -> str:
        """Extracts a readable batch name from folder string."""
        match = re.search(r'(?:RUN|R)[-_]?(\d+)', folder_name, re.IGNORECASE)
        if match:
            return f"R{int(match.group(1))}"

        for kw in self.config.control_keywords:
            if kw.lower() in folder_name.lower():
                return "Control"

        digits = re.findall(r'\d+', folder_name)
        if digits:
            return f"Batch_{digits[-1]}"
        return folder_name[:10]

    def _detect_user_initials(self, raw_df: pd.DataFrame) -> str:
        """Heuristic to guess user initials from filenames."""
        candidates = []
        if not raw_df.empty and 'SortKey' in raw_df.columns:
            # noinspection PySpellCheckingInspection
            for fname in raw_df['SortKey'].unique():
                match = re.search(r'-(?!RUN)([A-Z]{2,4})\d+', str(fname))
                if match:
                    candidates.append(match.group(1))

        if candidates:
            return Counter(candidates).most_common(1)[0][0]
        return self.config.default_initials

    def load_data(self, root_dir: Union[str, Path]) -> Tuple[pd.DataFrame, pd.DataFrame, str]:
        """Scans directory and returns raw data, batch mapping, and user initials."""
        root_path = Path(root_dir)
        logger.info(f"Scanning directory: {root_path.resolve()}")

        data_list = []
        batch_map_list = []
        target_files = []

        # 1. Scan files
        for dirpath, _, files in os.walk(root_path):
            if self.stop_requested:
                break

            matched_file = next(
                (f for f in files if any(fnmatch.fnmatch(f, p) for p in self.config.input_patterns)),
                None
            )

            if matched_file:
                target_files.append((Path(dirpath), matched_file))

        # Sort by folder name naturally
        target_files.sort(key=lambda x: self.natural_keys(str(x[0])))
        seen_labels: Dict[str, int] = {}

        # 2. Read files
        for dirpath, filename in target_files:
            if self.stop_requested:
                break

            file_path = dirpath / filename
            try:
                # Try UTF-8, fallback to GBK (common in lab equipment)
                try:
                    df_temp = pd.read_csv(file_path, encoding='utf-8', on_bad_lines='skip')
                except UnicodeDecodeError:
                    df_temp = pd.read_csv(file_path, encoding='gbk', on_bad_lines='skip')
                except PermissionError:
                    logger.error(f"Permission denied (File open?): {file_path}")
                    continue

                cols_map = self._normalize_columns(df_temp)
                if 'Eff' not in cols_map:
                    continue

                # Rename and subset
                rename_dict = {v: k for k, v in cols_map.items()}
                keep_cols = list(cols_map.values())
                if 'CellName' in cols_map:
                    keep_cols.append(cols_map['CellName'])

                df_subset = df_temp[keep_cols].rename(columns=rename_dict)
                df_subset = df_subset.loc[:, ~df_subset.columns.duplicated()]

                # Coerce numeric types
                for col in ANALYSIS_PARAMS:
                    if col in df_subset.columns:
                        df_subset[col] = pd.to_numeric(df_subset[col], errors='coerce')

                df_subset.dropna(subset=['Eff'], inplace=True)

                # Generate Batch Label
                short_label = self._parse_batch_label(dirpath.name)
                if short_label in seen_labels:
                    seen_labels[short_label] += 1
                    short_label = f"{short_label}-{seen_labels[short_label]}"
                else:
                    seen_labels[short_label] = 0

                df_subset['Batch'] = short_label
                df_subset['SortKey'] = str(dirpath)

                data_list.append(df_subset)
                batch_map_list.append({'Batch': short_label, 'Folder': dirpath.name})
                logger.info(f"Loaded: {short_label} ({len(df_subset)} cells)")

            except Exception as e:
                logger.warning(f"Error reading {file_path}: {e}")

        if not data_list:
            return pd.DataFrame(), pd.DataFrame(), self.config.default_initials

        raw_df = pd.concat(data_list, ignore_index=True)
        batch_map_df = pd.DataFrame(batch_map_list)
        user_initials = self._detect_user_initials(raw_df)

        return raw_df, batch_map_df, user_initials


# ================= MODULE 2: STATISTICS ENGINE =================

class IVStatistics:
    """Handles data cleaning, statistical analysis, and yield calculation."""

    def __init__(self, config: AnalyzerConfig):
        self.config = config

    def clean_data(self, raw_df: pd.DataFrame) -> Tuple[pd.DataFrame, List[str]]:
        """Applies filters and outlier removal."""
        if raw_df.empty:
            return pd.DataFrame(), []

        logger.info("Cleaning and processing data...")
        df = raw_df.copy()
        thresh = self.config.thresholds

        # 1. Filter by Scan Direction (Vectorized string op)
        if 'ScanDir' in df.columns:
            pref = self.config.scan_direction
            if pref in ['Reverse', 'Forward']:
                mask = df['ScanDir'].astype(str).str.contains(pref[0], case=False, na=False)
                df = df[mask]

        # 2. Physical Thresholds (Vectorized boolean mask)
        mask_phys = (
                (df['Eff'] > thresh['Eff_Min']) &
                (df['Voc'] > thresh['Voc_Min']) &
                (df['Jsc'] > thresh['Jsc_Min'])
        )
        if 'FF' in df.columns:
            mask_phys &= (df['FF'] > thresh['FF_Min']) & (df['FF'] < thresh['FF_Max'])

        df = df[mask_phys].copy()

        # 3. Remove Duplicates (Keep best Eff per cell)
        if self.config.remove_duplicates and 'CellName' in df.columns:
            df = df.sort_values('Eff', ascending=False).drop_duplicates(subset=['Batch', 'CellName'])

        # 4. Outlier Removal (IQR)
        cleaned_groups = []
        unique_batches = sorted(df['Batch'].unique(), key=IVDataLoader.natural_keys)

        for batch in unique_batches:
            group = df[df['Batch'] == batch]
            if self.config.outlier_removal and len(group) >= 5:
                q1 = group['Eff'].quantile(0.25)
                q3 = group['Eff'].quantile(0.75)
                iqr = q3 - q1
                group = group[(group['Eff'] >= q1 - 1.5 * iqr) & (group['Eff'] <= q3 + 3.0 * iqr)]
            cleaned_groups.append(group)

        if not cleaned_groups:
            logger.warning("All data removed during cleaning.")
            return pd.DataFrame(), []

        clean_df = pd.concat(cleaned_groups, ignore_index=True)
        clean_df['Eff_Bin'] = pd.cut(clean_df['Eff'], bins=EFF_BINS, labels=EFF_BIN_LABELS)

        return clean_df, unique_batches

    def assign_colors(self, batch_order: List[str]) -> Dict[str, str]:
        """Assigns consistent colors to batches."""
        palette = self.config.plotting.colors
        return {batch: palette[i % len(palette)] for i, batch in enumerate(batch_order)}

    def compute_statistics(self, clean_df: pd.DataFrame, batch_order: List[str]) -> Tuple[
        pd.DataFrame, pd.DataFrame, pd.DataFrame, pd.DataFrame, Dict[str, Any]]:
        """
        Calculates statistics using vectorized Pandas aggregation.
        """
        # --- OPTIMIZATION: Vectorized Aggregation ---
        available_params = [p for p in ANALYSIS_PARAMS if p in clean_df.columns]

        # Define aggregation dictionary
        agg_funcs = {param: ['mean', 'median', 'max', 'std'] for param in available_params}
        agg_funcs['Batch'] = ['count']  # Count samples

        # Perform aggregation
        grouped = clean_df.groupby('Batch')
        stats_raw = grouped.agg(agg_funcs)

        # Flatten MultiIndex columns (e.g., ('Eff', 'mean') -> 'Eff_Mean')
        stats_list = []

        # Identify Control Batch
        control_batch = next(
            (b for b in batch_order if any(k.lower() in b.lower() for k in self.config.control_keywords)),
            batch_order[0] if batch_order else None
        )
        comparisons = {'Control': control_batch, 'Results': {}}
        control_data = clean_df[clean_df['Batch'] == control_batch]

        # Reconstruct stats dataframe in correct order and perform T-Tests
        for batch in batch_order:
            if batch not in stats_raw.index:
                continue

            # Extract aggregated stats
            row = {'Batch': batch, 'Count': stats_raw.loc[batch, ('Batch', 'count')]}
            for param in available_params:
                row[f'{param}_Mean'] = stats_raw.loc[batch, (param, 'mean')]
                row[f'{param}_Median'] = stats_raw.loc[batch, (param, 'median')]
                row[f'{param}_Max'] = stats_raw.loc[batch, (param, 'max')]
                row[f'{param}_Std'] = stats_raw.loc[batch, (param, 'std')]

                # T-Test Logic (Cannot be easily vectorized as it compares groups)
                if batch != control_batch and not control_data.empty:
                    df_b = clean_df[clean_df['Batch'] == batch]
                    if len(df_b) > 1 and len(control_data) > 1:
                        try:
                            a = df_b[param].dropna()
                            b = control_data[param].dropna()
                            t, p = stats.ttest_ind(a, b, equal_var=False)
                            diff = row[f'{param}_Mean'] - control_data[param].mean()
                            sig = '***' if p < 0.001 else '**' if p < 0.01 else '*' if p < 0.05 else 'ns'

                            if batch not in comparisons['Results']:
                                comparisons['Results'][batch] = {}
                            comparisons['Results'][batch][param] = {
                                'p': p, 'sig': sig, 'diff': diff,
                                'dir': "Increase" if diff > 0 else "Decrease"
                            }
                        except Exception as e:
                            logger.debug(f"T-Test skipped for {batch} vs {control_batch} on {param}: {e}")

            stats_list.append(row)

        stats_df = pd.DataFrame(stats_list)

        # Champion Selection (Vectorized)
        champion_df = clean_df.sort_values('Eff', ascending=False).groupby('Batch', as_index=False).first()
        champion_df = champion_df.set_index('Batch').reindex(batch_order).reset_index()

        # Top 10 Cells
        top_cells_df = clean_df.sort_values('Eff', ascending=False).head(10)

        # Yield Calculation
        yield_raw = pd.crosstab(clean_df['Batch'], clean_df['Eff_Bin'], normalize='index') * 100
        yield_df = yield_raw.reindex(batch_order).fillna(0).reset_index()

        return stats_df, champion_df, top_cells_df, yield_df, comparisons


# ================= MODULE 3: VISUALIZER =================

class IVVisualizer:
    """Generates plots using Matplotlib/Seaborn."""

    def __init__(self, config: AnalyzerConfig, output_dir: Optional[Path]):
        self.config = config
        self.output_dir = output_dir
        self.img_paths: Dict[str, Path] = {}

        # State injection (Facade pattern trade-off)
        self.clean_df: Optional[pd.DataFrame] = None
        self.stats_df: Optional[pd.DataFrame] = None
        self.champion_df: Optional[pd.DataFrame] = None
        self.batch_order: List[str] = []
        self.group_colors: Dict[str, str] = {}

        self._setup_plot_style()

    def _setup_plot_style(self) -> None:
        sns.set_theme(style="ticks", context="talk", font_scale=1.1)
        # noinspection SpellCheckingInspection
        plt.rcParams.update({
            'font.family': 'sans-serif',
            'font.sans-serif': self.config.plotting.font_family,
            'axes.unicode_minus': False,
            'lines.linewidth': self.config.plotting.line_width,
            'figure.facecolor': 'white',
            'axes.grid': True,
            'grid.alpha': 0.3,
            'grid.linestyle': '--',
            'xtick.direction': 'in', 'ytick.direction': 'in', 'axes.linewidth': 1.5
        })

    def _save_plot(self, filename: str) -> Optional[Path]:
        if not self.output_dir:
            return None
        try:
            path = self.output_dir / filename
            plt.savefig(path, dpi=self.config.plotting.dpi, bbox_inches='tight')
            return path
        except Exception as e:
            logger.error(f"Failed to save {filename}: {e}")
            return None
        finally:
            plt.close()

    @staticmethod
    def _apply_legend_style(ax):
        ax.legend(bbox_to_anchor=(1.05, 1), loc='upper left', borderaxespad=0., frameon=True)

    @staticmethod
    def _reconstruct_jv_curve(voc, jsc, ff, points=100):
        """Reconstructs JV curve using single-diode model approximation."""
        # Fix: Cast to float to avoid numpy division ambiguity warnings
        voc_val = float(voc)
        jsc_val = float(jsc)

        v_curve = np.linspace(-0.1, voc_val + 0.05, points)

        # Default linear fallback
        j_curve = jsc_val * (1 - v_curve / voc_val)

        try:
            # Fix: Rename argument 'B' to 'diode_factor' for PEP 8
            def model_j(v, diode_factor):
                with np.errstate(over='ignore', invalid='ignore'):
                    # Fix: Ensure types are float for division
                    term = np.exp(voc_val / diode_factor) - 1
                    if term == 0 or np.isinf(term) or np.isnan(term):
                        term = 1e9
                    exp_v = np.exp(v / diode_factor) - 1
                    return jsc_val - (jsc_val / term) * exp_v

            # Fix: Rename 'B' to 'diode_factor'
            def get_ff_error(diode_factor):
                if diode_factor <= 1e-4: return 100
                v_r = np.linspace(0, voc_val, 50)
                j_r = model_j(v_r, diode_factor)
                p_max = np.max(v_r * j_r)
                return (p_max / (voc_val * jsc_val)) * 100 - ff

            b_opt = optimize.fsolve(get_ff_error, 0.05, xtol=1e-4)[0]
            if b_opt < 1e-4:
                b_opt = 0.05

            j_curve = model_j(v_curve, b_opt)
        except Exception as e:
            logger.debug(f"JV model fit failed (using linear fallback): {e}")

        return v_curve, j_curve

    # --- PLOTTING METHODS ---

    def _plot_boxplot(self) -> None:
        try:
            params = [p for p in ANALYSIS_PARAMS if p in self.clean_df.columns]
            fig, axes = plt.subplots(1, len(params), figsize=(5 * len(params), 6))
            if len(params) == 1:
                axes = [axes]

            for ax, param in zip(axes, params):
                sns.boxplot(
                    data=self.clean_df, x='Batch', y=param, order=self.batch_order,
                    ax=ax, palette=self.group_colors, showfliers=False, width=0.5
                )
                if self.config.plotting.show_points:
                    sns.stripplot(
                        data=self.clean_df, x='Batch', y=param, order=self.batch_order,
                        ax=ax, color=".25", size=4, alpha=0.6, jitter=True
                    )
                ax.set_title("")
                ax.set_ylabel(f"{param} {UNIT_MAP.get(param, '')}", fontweight='bold')
                ax.tick_params(axis='x', rotation=45)
                ax.set_xlabel('')

            plt.tight_layout()
            self.img_paths['box'] = self._save_plot('1_Boxplot.png')
        except Exception as e:
            logger.error(f"Boxplot failed: {e}")

    def _plot_histogram(self) -> None:
        try:
            plt.figure(figsize=self.config.plotting.base_figsize)
            sns.histplot(
                data=self.clean_df, x='Eff', hue='Batch', hue_order=self.batch_order,
                palette=self.group_colors, kde=True, element="step", fill=False, legend=False
            )
            plt.title('Efficiency Distribution', fontweight='bold')
            plt.xlabel(f'Efficiency {UNIT_MAP["Eff"]}')

            handles = [
                mlines.Line2D([], [], color=self.group_colors[b], linewidth=2, label=b)
                for b in self.batch_order
            ]
            plt.legend(handles=handles, bbox_to_anchor=(1.05, 1), loc='upper left', title="Batch")

            self.img_paths['hist'] = self._save_plot('1_Histogram.png')
        except Exception as e:
            logger.error(f"Histogram failed: {e}")

    def _plot_trend(self) -> None:
        try:
            params = [p for p in ANALYSIS_PARAMS if p in self.clean_df.columns]
            fig, axes = plt.subplots(1, len(params), figsize=(5 * len(params), 6))
            if len(params) == 1:
                axes = [axes]

            for ax, param in zip(axes, params):
                td = self.stats_df.set_index('Batch').reindex(self.batch_order)
                x = range(len(self.batch_order))
                ax.plot(x, td[f'{param}_Max'], label='Max', color='#E64B35', linestyle='--', marker='^', markersize=8)
                ax.plot(x, td[f'{param}_Mean'], label='Mean', color='#3C5488', linestyle='-', marker='o', markersize=8)
                ax.plot(x, td[f'{param}_Median'], label='Median', color='#00A087', linestyle=':', marker='x',
                        markersize=8)
                ax.set_title("")
                ax.set_ylabel(f"{param} {UNIT_MAP.get(param, '')}", fontweight='bold')
                ax.set_xticks(x)
                ax.set_xticklabels(self.batch_order, rotation=45)

            handles, labels = axes[0].get_legend_handles_labels()
            axes[-1].legend(handles, labels, bbox_to_anchor=(1.05, 1), loc='upper left', borderaxespad=0.)

            plt.tight_layout()

            self.img_paths['trend'] = self._save_plot('2_Trend.png')
        except Exception as e:
            logger.error(f"Trend plot failed: {e}")

    def _plot_yield(self) -> None:
        try:
            bin_counts = pd.crosstab(self.clean_df['Batch'], self.clean_df['Eff_Bin'], normalize='index') * 100
            bin_counts = bin_counts.reindex(self.batch_order).fillna(0)
            bin_counts = bin_counts.loc[:, (bin_counts != 0).any(axis=0)]

            if not bin_counts.empty:
                ax = bin_counts.plot(
                    kind='bar', stacked=True, figsize=self.config.plotting.base_figsize,
                    colormap='viridis_r', width=0.7
                )
                plt.title('Efficiency Yield Distribution', fontweight='bold')
                plt.ylabel('Percentage (%)')
                plt.xticks(rotation=45)
                ax.legend(bbox_to_anchor=(1.05, 1), loc='upper left', title="Eff Range")
                self.img_paths['yield'] = self._save_plot('3_Yield.png')
        except Exception as e:
            logger.error(f"Yield plot failed: {e}")

    def _plot_jv_curves(self) -> None:
        try:
            plt.figure(figsize=self.config.plotting.base_figsize)
            has_jv = False
            for batch in self.batch_order:
                if batch in self.champion_df['Batch'].values:
                    row = self.champion_df[self.champion_df['Batch'] == batch].iloc[0]
                    try:
                        v, j = self._reconstruct_jv_curve(row['Voc'], row['Jsc'], row['FF'])
                        plt.plot(
                            v, j, label=f"{batch} ({row['Eff']:.2f}%)",
                            color=self.group_colors.get(batch, 'black'), linewidth=2.5
                        )
                        has_jv = True
                    except Exception as e:
                        logger.warning(f"Could not plot JV curve for {batch}: {e}")
                        continue
            if has_jv:
                plt.axhline(0, color='gray', linestyle='--', linewidth=0.8)
                plt.axvline(0, color='gray', linestyle='--', linewidth=0.8)
                plt.xlabel('Voltage (V)', fontweight='bold')
                plt.ylabel('Current Density (mA/cm²)', fontweight='bold')
                plt.title('J-V Curves of Champion Cells', fontweight='bold')
                self._apply_legend_style(plt.gca())
                self.img_paths['jv_curve'] = self._save_plot('4_JV_Curves.png')
            else:
                plt.close()
        except Exception as e:
            logger.error(f"JV Curve plot failed: {e}")

    def _plot_correlations(self) -> None:
        try:
            if 'Voc' in self.clean_df.columns and 'Jsc' in self.clean_df.columns:
                plt.figure(figsize=self.config.plotting.base_figsize)
                sns.scatterplot(
                    data=self.clean_df, x='Jsc', y='Voc', hue='Batch', hue_order=self.batch_order,
                    palette=self.group_colors, s=self.config.plotting.marker_size, alpha=0.8, edgecolor='k'
                )

                avg_ff = self.clean_df['FF'].mean() / 100.0 if 'FF' in self.clean_df.columns else 0.75
                x_min, x_max = plt.xlim()
                y_min, y_max = plt.ylim()

                if x_max > x_min and y_max > y_min:
                    xx = np.linspace(x_min, x_max, 100)
                    yy = np.linspace(y_min, y_max, 100)

                    # Fix: Renamed meshgrid vars to lowercase to satisfy PEP 8 and local var warning
                    grid_x, grid_y = np.meshgrid(xx, yy)
                    grid_z = grid_x * grid_y * avg_ff

                    cs = plt.contour(grid_x, grid_y, grid_z, levels=8, colors='gray', linestyles='--', alpha=0.4)
                    plt.clabel(cs, inline=1, fontsize=10, fmt='%.1f%%')

                plt.title(f'Voc vs. Jsc Correlation (Iso-Eff @ FF={avg_ff * 100:.0f}%)', fontweight='bold')
                self._apply_legend_style(plt.gca())
                self.img_paths['voc_jsc'] = self._save_plot('5_Voc_Jsc_Tradeoff.png')

            fig, axes = plt.subplots(2, 2, figsize=(12, 10))
            pairs = [
                ('Voc', 'Eff', axes[0, 0]), ('Jsc', 'Eff', axes[0, 1]),
                ('FF', 'Eff', axes[1, 0]), ('Voc', 'Jsc', axes[1, 1])
            ]

            for x, y, ax in pairs:
                if x in self.clean_df.columns and y in self.clean_df.columns:
                    sns.scatterplot(
                        data=self.clean_df, x=x, y=y, hue='Batch', hue_order=self.batch_order,
                        palette=self.group_colors, s=self.config.plotting.marker_size, alpha=0.8, ax=ax,
                        legend=False
                    )
                    ax.set_title(f'{y} vs {x}', fontweight='bold')
                    ax.set_xlabel(f'{x} {UNIT_MAP.get(x, "")}')
                    ax.set_ylabel(f'{y} {UNIT_MAP.get(y, "")}')

            handles = [
                plt.Line2D([0], [0], marker='o', color='w', markerfacecolor=self.group_colors[b], markersize=10)
                for b in self.batch_order
            ]
            fig.legend(handles, self.batch_order, bbox_to_anchor=(0.86, 0.9), loc='upper left', title="Batch")
            plt.tight_layout(rect=(0, 0, 0.85, 1))
            self.img_paths['combo_drivers'] = self._save_plot('6_Drivers.png')
        except Exception as e:
            logger.error(f"Correlation plots failed: {e}")

    def _plot_resistance_analysis(self) -> None:
        try:
            if 'Rs' not in self.clean_df.columns or 'FF' not in self.clean_df.columns: return
            fig, axes = plt.subplots(2, 2, figsize=(12, 10))

            # FF vs Rs
            sns.scatterplot(
                data=self.clean_df, x='Rs', y='FF', hue='Batch', hue_order=self.batch_order,
                palette=self.group_colors, s=self.config.plotting.marker_size, alpha=0.8, ax=axes[0, 0],
                legend=False
            )
            axes[0, 0].set_title('FF vs Rs', fontweight='bold')
            axes[0, 0].set_xlabel(f'Rs {UNIT_MAP["Rs"]}')

            # FF vs Rsh
            if 'Rsh' in self.clean_df.columns:
                sns.scatterplot(
                    data=self.clean_df, x='Rsh', y='FF', hue='Batch', hue_order=self.batch_order,
                    palette=self.group_colors, s=self.config.plotting.marker_size, alpha=0.8, ax=axes[0, 1],
                    legend=False
                )
                axes[0, 1].set_title('FF vs Rsh', fontweight='bold')
                axes[0, 1].set_xlabel(f'Rsh {UNIT_MAP["Rsh"]}')
                axes[0, 1].set_xscale('log')

            # Eff vs Rs
            sns.scatterplot(
                data=self.clean_df, x='Rs', y='Eff', hue='Batch', hue_order=self.batch_order,
                palette=self.group_colors, s=self.config.plotting.marker_size, alpha=0.8, ax=axes[1, 0],
                legend=False
            )
            axes[1, 0].set_title('Eff vs Rs', fontweight='bold')
            axes[1, 0].set_xlabel(f'Rs {UNIT_MAP["Rs"]}')

            # Eff vs Rsh
            if 'Rsh' in self.clean_df.columns:
                sns.scatterplot(
                    data=self.clean_df, x='Rsh', y='Eff', hue='Batch', hue_order=self.batch_order,
                    palette=self.group_colors, s=self.config.plotting.marker_size, alpha=0.8, ax=axes[1, 1],
                    legend=False
                )
                axes[1, 1].set_title('Eff vs Rsh', fontweight='bold')
                axes[1, 1].set_xlabel(f'Rsh {UNIT_MAP["Rsh"]}')
                axes[1, 1].set_xscale('log')

            handles = [
                plt.Line2D([0], [0], marker='o', color='w', markerfacecolor=self.group_colors[b], markersize=10)
                for b in self.batch_order
            ]
            fig.legend(handles, self.batch_order, bbox_to_anchor=(0.86, 0.9), loc='upper left', title="Batch")
            plt.tight_layout(rect=(0, 0, 0.85, 1))
            self.img_paths['resistance'] = self._save_plot('7_Resistance.png')
        except Exception as e:
            logger.error(f"Resistance plot failed: {e}")

    def visualize(self, clean_df: pd.DataFrame, stats_df: pd.DataFrame,
                  champion_df: pd.DataFrame, batch_order: List[str],
                  group_colors: Dict[str, str]) -> Dict[str, Path]:
        """Main entry point for visualization."""
        if clean_df.empty:
            return {}

        # Inject state
        self.clean_df = clean_df
        self.stats_df = stats_df
        self.champion_df = champion_df
        self.batch_order = batch_order
        self.group_colors = group_colors

        logger.info("Generating plots...")

        try:
            self._plot_boxplot()
            self._plot_histogram()
            self._plot_trend()
            self._plot_yield()
            self._plot_jv_curves()
            self._plot_correlations()
            self._plot_resistance_analysis()
        except Exception as e:
            logger.error(f"Critical error during visualization: {e}", exc_info=True)
        finally:
            plt.close('all')
            self.clean_df = None
            self.stats_df = None

        return self.img_paths


# ================= MODULE 4: REPORT GENERATOR =================

class IVReportGenerator:
    """Generates Excel, Word, and PowerPoint reports."""

    def __init__(self, config: AnalyzerConfig):
        self.config = config
        self.output_dir: Optional[Path] = None

    def create_output_dir(self, root_dir: Path, raw_df: pd.DataFrame, user_initials: str) -> Path:
        """Creates output directory with timestamp and run info."""
        run_nums = []
        if not raw_df.empty and 'Batch' in raw_df.columns:
            for batch in raw_df['Batch'].unique():
                match = re.search(r'\d+', str(batch))
                if match:
                    run_nums.append(int(match.group(0)))

        suffix = datetime.now().strftime("%Y%m%d_%H%M")
        if run_nums:
            min_r, max_r = min(run_nums), max(run_nums)
            suffix = f"RUN{min_r}" if min_r == max_r else f"RUN{min_r}-{max_r}"

        folder_name = f"IV_Report_{user_initials}_{suffix}"
        self.output_dir = root_dir / folder_name
        self.output_dir.mkdir(exist_ok=True)
        logger.info(f"Output directory created: {self.output_dir}")
        return self.output_dir

    def export_reports(self, clean_df: pd.DataFrame, stats_df: pd.DataFrame,
                       champion_df: pd.DataFrame, top_cells_df: pd.DataFrame,
                       yield_df: pd.DataFrame, batch_map_df: pd.DataFrame,
                       comparisons: Dict[str, Any], img_paths: Dict[str, Path],
                       user_initials: str) -> None:

        if not self.output_dir:
            return
        logger.info("Exporting reports...")

        try:
            # --- 1. Excel Export ---
            path = self.output_dir / self.config.excel_data_name
            with pd.ExcelWriter(path, engine='openpyxl') as writer:
                clean_df.to_excel(writer, sheet_name='Cleaned_Data', index=False)
                stats_df.to_excel(writer, sheet_name='Statistics', index=False)
                if not champion_df.empty:
                    champion_df.to_excel(writer, sheet_name='Champions', index=False)
                top_cells_df.to_excel(writer, sheet_name='Top_10_Cells', index=False)
                yield_df.to_excel(writer, sheet_name='Yield_Table', index=False)

                for sheet in writer.sheets.values():
                    for col in sheet.columns:
                        sheet.column_dimensions[col[0].column_letter].width = 15

            # --- 2. Word Export ---
            doc = Document()
            doc.add_heading(f'IV Analysis Report ({user_initials})', 0)
            doc.add_heading('1. Executive Summary', level=1)

            if comparisons.get('Results'):
                ctrl = comparisons['Control']
                for batch, res in comparisons['Results'].items():
                    p = doc.add_paragraph()
                    run = p.add_run(f"▶ {batch} vs {ctrl}: ")
                    run.bold = True
                    run.font.color.rgb = RGBColor(0, 51, 102)
                    if 'Eff' in res:
                        e = res['Eff']
                        p.add_run(f"Eff {e['dir']} by {abs(e['diff']):.2f}% (p={e['p']:.3f}, {e['sig']}).")
            else:
                doc.add_paragraph("Single batch analysis.")

            doc.add_heading('2. Data Tables', level=1)
            self._add_word_table(doc, batch_map_df, ['Batch', 'Folder'], "2.1 Batch & Folder Reference")

            cols = ['Batch', 'Count', 'Eff_Mean', 'Eff_Max', 'Voc_Mean', 'FF_Mean']
            cols = [c for c in cols if c in stats_df.columns]
            self._add_word_table(doc, stats_df, cols, "2.2 Statistical Summary of PV Parameters")

            champ_cols = [c for c in ['Batch', 'CellName', 'Eff', 'Voc', 'Jsc', 'FF', 'Rs'] if c in champion_df.columns]
            self._add_word_table(doc, champion_df, champ_cols, "2.3 Champion Cell Parameters")
            self._add_word_table(doc, top_cells_df, champ_cols, "2.4 Top 10 Highest Efficiency Cells")

            yield_cols = [c for c in yield_df.columns if c != 'Batch']
            valid_yield_cols = ['Batch'] + [c for c in yield_cols if yield_df[c].sum() > 0]
            self._add_word_table(doc, yield_df, valid_yield_cols, "2.5 Efficiency Yield Distribution (%)")

            doc.add_heading('3. Visual Analysis', level=1)
            plot_order = [
                ('jv_curve', 'J-V Curves of Champion Cells'),
                ('voc_jsc', 'Voc vs. Jsc Correlation Analysis'),
                ('box', 'Distribution of Electrical Parameters'),
                ('hist', 'Efficiency Histogram'),
                ('trend', 'Batch Trend Analysis (Max/Mean/Median)'),
                ('yield', 'Efficiency Yield Stack'),
                ('combo_drivers', 'Key Efficiency Drivers (Correlations)'),
                ('resistance', 'Parasitic Resistance Analysis (Rs & Rsh)')
            ]
            for key, title in plot_order:
                if key in img_paths and img_paths[key]:
                    doc.add_heading(title, level=2)
                    # Fix: Explicit cast to str to satisfy PyCharm Type Check (Path vs str)
                    doc.add_picture(str(img_paths[key]), width=Inches(6.5))
            doc.save(self.output_dir / self.config.report_docx_name)

            # --- 3. PPT Export ---
            if HAS_PPTX:
                prs = Presentation()
                prs.slide_width = PptInches(13.33)
                prs.slide_height = PptInches(7.5)

                # Title Slide
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = "IV Measurement Analysis"
                slide.placeholders[1].text = f"Operator: {user_initials}\nDate: {datetime.now().strftime('%Y-%m-%d')}"

                # Summary Slide
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Executive Summary"
                tf = slide.shapes.placeholders[1].text_frame
                slide.shapes.title.text_frame.paragraphs[0].font.size = PptPt(24)

                if comparisons.get('Results'):
                    for batch, res in comparisons['Results'].items():
                        p = tf.add_paragraph()
                        p.text = f"{batch} vs {comparisons['Control']}:"
                        p.font.bold = True
                        p.font.size = PptPt(16)

                        if 'Eff' in res:
                            p2 = tf.add_paragraph()
                            p2.level = 1
                            p2.text = f"Efficiency {res['Eff']['dir']} {abs(res['Eff']['diff']):.2f}% ({res['Eff']['sig']})"
                            p2.font.size = PptPt(16)

                self._add_pptx_table(prs, batch_map_df, ['Batch', 'Folder'], "Batch & Folder Reference")
                self._add_pptx_table(prs, stats_df, cols, "Statistical Summary of PV Parameters")
                self._add_pptx_table(prs, champion_df, champ_cols, "Champion Cell Parameters")
                self._add_pptx_table(prs, top_cells_df, champ_cols, "Top 10 Highest Efficiency Cells")
                self._add_pptx_table(prs, yield_df, valid_yield_cols, "Efficiency Yield Distribution (%)")

                for key, title in plot_order:
                    if key in img_paths and img_paths[key]:
                        slide = prs.slides.add_slide(prs.slide_layouts[5])
                        slide.shapes.title.text = title
                        slide.shapes.title.text_frame.paragraphs[0].font.size = PptPt(24)
                        # Fix: Explicit cast to str to satisfy PyCharm Type Check
                        slide.shapes.add_picture(str(img_paths[key]), PptInches(2.5), PptInches(1.2),
                                                 height=PptInches(5.8))

                prs.save(self.output_dir / self.config.report_pptx_name)

            logger.info(f"Reports saved: {self.output_dir}")
        except Exception as e:
            logger.error(f"Report generation failed: {e}", exc_info=True)

    # --- TABLE GENERATION METHODS ---

    @staticmethod
    def _add_word_table(doc, df, columns, title):
        doc.add_heading(title, level=2)
        if df.empty or not columns: return

        table = doc.add_table(rows=1, cols=len(columns))
        table.style = 'Table Grid'
        hdr_cells = table.rows[0].cells

        for i, col in enumerate(columns):
            header_text = str(col).replace('_', ' (') + ')' if '_' in str(col) else str(col)
            hdr_cells[i].text = header_text
            p = hdr_cells[i].paragraphs[0]
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.runs[0]
            run.font.bold = True
            run.font.size = Pt(9)
            hdr_cells[i].vertical_alignment = WD_ALIGN_VERTICAL.CENTER

        for _, row in df.iterrows():
            row_cells = table.add_row().cells
            for i, col in enumerate(columns):
                val = row.get(col, "")
                if isinstance(val, (int, float)):
                    txt = f"{val:.0f}" if 'Count' in col else f"{val:.2f}"
                else:
                    txt = str(val)

                row_cells[i].text = txt
                p = row_cells[i].paragraphs[0]
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                p.runs[0].font.size = Pt(9)
                row_cells[i].vertical_alignment = WD_ALIGN_VERTICAL.CENTER
        doc.add_paragraph()

    @staticmethod
    def _add_pptx_table(prs, df, columns, title):
        if df.empty or not columns: return
        max_rows = 10
        total_rows = len(df)
        num_slides = math.ceil(total_rows / max_rows)

        for i in range(num_slides):
            start = i * max_rows
            end = min((i + 1) * max_rows, total_rows)
            chunk = df.iloc[start:end]

            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"{title} ({i + 1}/{num_slides})" if num_slides > 1 else title
            slide.shapes.title.text_frame.paragraphs[0].font.size = PptPt(24)

            rows, cols = len(chunk) + 1, len(columns)
            left, top, width, height = PptInches(0.5), PptInches(1.5), PptInches(12.0), PptInches(0.5 * rows)

            table = slide.shapes.add_table(rows, cols, left, top, width, height).table

            for c_idx, col_name in enumerate(columns):
                cell = table.cell(0, c_idx)
                header_text = str(col_name).replace('_', ' (') + ')' if '_' in str(col_name) else str(col_name)
                cell.text = header_text
                p = cell.text_frame.paragraphs[0]
                p.font.bold = True
                p.font.size = PptPt(12)
                p.alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for r_idx in range(len(chunk)):
                row_data = chunk.iloc[r_idx]
                for c_idx, col_name in enumerate(columns):
                    cell = table.cell(r_idx + 1, c_idx)
                    val = row_data.get(col_name, "")
                    if isinstance(val, (int, float)):
                        txt = f"{val:.0f}" if 'Count' in col_name else f"{val:.2f}"
                    else:
                        txt = str(val)

                    cell.text = txt
                    p = cell.text_frame.paragraphs[0]
                    p.font.size = PptPt(11)
                    p.alignment = PP_ALIGN.CENTER
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


# ================= FACADE: MAIN CONTROLLER =================

class IVBatchAnalyzer:
    """Facade class that orchestrates the modular components."""

    def __init__(self, root_dir: Union[str, Path], config: Optional[AnalyzerConfig] = None):
        self.root_dir = Path(root_dir)
        self.config = config or AnalyzerConfig()
        self.stop_requested = False

        self.loader = IVDataLoader(self.config)
        self.stats = IVStatistics(self.config)
        self.reporter = IVReportGenerator(self.config)

    def run(self) -> None:
        # 1. Load Data
        self.loader.stop_requested = self.stop_requested
        raw_df, batch_map_df, user_initials = self.loader.load_data(self.root_dir)

        if self.stop_requested:
            logger.info("Analysis cancelled by user.")
            return

        if raw_df.empty:
            logger.error("No data found or operation cancelled.")
            return

        # 2. Create Output Directory
        output_dir = self.reporter.create_output_dir(self.root_dir, raw_df, user_initials)

        # 3. Process Statistics
        if self.stop_requested: return
        clean_df, batch_order = self.stats.clean_data(raw_df)

        if clean_df.empty:
            logger.warning("No data remained after cleaning.")
            return

        group_colors = self.stats.assign_colors(batch_order)
        stats_df, champion_df, top_cells_df, yield_df, comparisons = \
            self.stats.compute_statistics(clean_df, batch_order)

        # 4. Visualize
        if self.stop_requested: return
        visualizer = IVVisualizer(self.config, output_dir)
        img_paths = visualizer.visualize(clean_df, stats_df, champion_df, batch_order, group_colors)

        # 5. Generate Reports
        if self.stop_requested: return
        self.reporter.export_reports(
            clean_df, stats_df, champion_df, top_cells_df, yield_df,
            batch_map_df, comparisons, img_paths, user_initials
        )

        logger.info("=== Done ===")


# ================= GUI =================

class TextHandler(logging.Handler):
    """Thread-safe logging handler for Tkinter Text widget."""

    def __init__(self, text_widget):
        super().__init__()
        self.text_widget = text_widget

    def emit(self, record):
        msg = self.format(record)
        self.text_widget.after(0, self._append_text, msg)

    def _append_text(self, msg):
        try:
            self.text_widget.configure(state='normal')
            self.text_widget.insert(tk.END, msg + '\n')
            self.text_widget.see(tk.END)
            self.text_widget.configure(state='disabled')
        except Exception:
            pass


class IVAnalyzerGUI:
    def __init__(self, root):
        self.root = root
        self.root.title("IV Batch Analyzer V4.2 (Modular)")
        self.root.geometry("750x650")

        self.folder_path = tk.StringVar(value=str(Path.cwd()))
        self.scan_dir_var = tk.StringVar(value="Reverse")
        self.eff_min_var = tk.DoubleVar(value=0.1)
        self.remove_outliers_var = tk.BooleanVar(value=True)
        self.status_var = tk.StringVar(value="Ready")

        self.analyzer_instance = None
        self._setup_ui()
        self._setup_logging()

    def _setup_ui(self):
        style = ttk.Style()
        style.theme_use('clam')

        # Config Frame
        cf = ttk.LabelFrame(self.root, text="Configuration", padding=10)
        cf.pack(fill="x", padx=10, pady=5)

        ttk.Label(cf, text="Folder:").grid(row=0, column=0, sticky="w")
        ttk.Entry(cf, textvariable=self.folder_path, width=55).grid(row=0, column=1, padx=5)
        ttk.Button(cf, text="Browse", command=self.select_folder).grid(row=0, column=2)

        sf = ttk.Frame(cf)
        sf.grid(row=1, column=0, columnspan=3, sticky="w", pady=10)

        ttk.Label(sf, text="Scan:").pack(side="left")
        ttk.Combobox(
            sf, textvariable=self.scan_dir_var, values=["Reverse", "Forward", "All"],
            state="readonly", width=8
        ).pack(side="left", padx=5)

        ttk.Label(sf, text="Min Eff:").pack(side="left", padx=(15, 0))
        ttk.Entry(sf, textvariable=self.eff_min_var, width=5).pack(side="left", padx=5)
        ttk.Checkbutton(sf, text="Remove Outliers", variable=self.remove_outliers_var).pack(side="left", padx=15)

        # Log Frame
        lf = ttk.LabelFrame(self.root, text="Log", padding=10)
        lf.pack(fill="both", expand=True, padx=10, pady=5)
        self.log_text = scrolledtext.ScrolledText(lf, state='disabled', height=15, font=("Consolas", 9))
        self.log_text.pack(fill="both", expand=True)

        # Action Frame
        af = ttk.Frame(self.root, padding=10)
        af.pack(fill="x")
        self.run_btn = ttk.Button(af, text="▶ Run Analysis", command=self.start_thread)
        self.run_btn.pack(side="right", padx=5)

        # Status Bar
        # Fix: Use string literal "sunken" instead of tk.SUNKEN for stricter type compliance.
        ttk.Label(self.root, textvariable=self.status_var, relief="sunken", anchor="w").pack(side="bottom", fill="x")

    def _setup_logging(self):
        logger.addHandler(TextHandler(self.log_text))

    def select_folder(self):
        p = filedialog.askdirectory()
        if p:
            self.folder_path.set(p)

    def start_thread(self):
        self.run_btn.config(state="disabled")
        self.status_var.set("Running...")
        self.log_text.configure(state='normal')
        self.log_text.delete(1.0, tk.END)
        self.log_text.configure(state='disabled')
        threading.Thread(target=self._run, daemon=True).start()

    def _run(self):
        try:
            cfg = AnalyzerConfig(
                scan_direction=self.scan_dir_var.get(),
                outlier_removal=self.remove_outliers_var.get(),
                thresholds={
                    "Eff_Min": self.eff_min_var.get(),
                    "Voc_Min": 0.1, "Jsc_Min": 0.1, "FF_Min": 10.0, "FF_Max": 90.0
                }
            )

            self.analyzer_instance = IVBatchAnalyzer(self.folder_path.get(), cfg)
            self.analyzer_instance.run()

            self.root.after(0, lambda: messagebox.showinfo("Success", "Analysis Complete!"))
        # Fix: Suppress 'Too broad exception' warning as this is the top-level GUI thread wrapper
        # noinspection PyBroadException
        except Exception as e:
            logger.error(f"Error: {e}", exc_info=True)
            self.root.after(0, lambda: messagebox.showerror("Error", str(e)))
        finally:
            self.root.after(0, lambda: [self.run_btn.config(state="normal"), self.status_var.set("Ready")])


if __name__ == "__main__":
    # Renamed 'root' to 'main_window' to avoid shadowing the parameter in IVAnalyzerGUI.__init__
    main_window = tk.Tk()
    IVAnalyzerGUI(main_window)
    main_window.mainloop()