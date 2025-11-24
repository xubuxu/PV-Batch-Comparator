"""
IV Batch Analyzer V5.0 - Professional Edition
Reporter Module

Generates Excel, Word, and PowerPoint reports.
"""
from __future__ import annotations

import logging
import math
import re
from datetime import datetime
from pathlib import Path
from typing import Dict, Any, Optional

import pandas as pd
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_ALIGN_VERTICAL

try:
    from pptx import Presentation
    from pptx.util import Inches as PptInches, Pt as PptPt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    Presentation = None
    PptInches = None
    PptPt = None
    PP_ALIGN = None
    MSO_ANCHOR = None

logger = logging.getLogger(__name__)


# ================= CUSTOM EXCEPTIONS =================

class ReportGenerationError(Exception):
    """Raised when report generation fails."""
    pass


# ================= REPORT GENERATOR =================

class IVReportGenerator:
    """Generates Excel, Word, and PowerPoint reports."""

    def __init__(self, config):
        """
        Initialize IVReportGenerator.
        
        Args:
            config: AnalyzerConfig instance
        """
        self.config = config
        self.output_dir: Optional[Path] = None

    def create_output_dir(self, root_dir: Path, raw_df: pd.DataFrame, user_initials: str) -> Path:
        """
        Creates output directory with configurable naming convention.
        
        Args:
            root_dir: Root directory for analysis
            raw_df: Raw dataframe
            user_initials: User initials
            
        Returns:
            Path to created output directory
        """
        from .config import NAMING_PATTERNS
        
        # 1. Prepare context variables
        now = datetime.now()
        run_nums = []
        batch_names = []
        
        if not raw_df.empty and 'Batch' in raw_df.columns:
            batch_names = sorted(raw_df['Batch'].unique().astype(str))
            for batch in batch_names:
                match = re.search(r'\d+', str(batch))
                if match:
                    run_nums.append(int(match.group(0)))
        
        # Determine RunID
        if run_nums:
            min_r, max_r = min(run_nums), max(run_nums)
            run_id = f"RUN{min_r}" if min_r == max_r else f"RUN{min_r}-{max_r}"
        else:
            run_id = "RUN001"
            
        # Determine Batch Summary
        if len(batch_names) == 1:
            batch_summary = batch_names[0]
        elif len(batch_names) > 1:
            batch_summary = f"{batch_names[0]}_etal"
        else:
            batch_summary = "BatchData"
            
        context = {
            "User": user_initials,
            "Timestamp": now.strftime("%Y%m%d_%H%M"),
            "Timestamp_Short": now.strftime("%H%M%S"),
            "YYYYMMDD": now.strftime("%Y%m%d"),
            "HHMM": now.strftime("%H%M"),
            "Batch_Summary": batch_summary,
            "RunID": run_id
        }
        
        # 2. Get pattern
        convention = getattr(self.config, 'naming_convention', 'Traceable')
        pattern = NAMING_PATTERNS.get(convention, NAMING_PATTERNS['Traceable'])
        
        # 3. Format folder name
        try:
            folder_name = pattern.format(**context)
        except KeyError as e:
            logger.warning(f"Invalid naming pattern key: {e}. Falling back to default.")
            folder_name = f"IV_Report_{user_initials}_{context['Timestamp']}"
            
        # Sanitize (remove illegal chars for Windows/Linux)
        folder_name = re.sub(r'[<>:"/\\|?*]', '_', folder_name)
        
        self.output_dir = root_dir / folder_name
        self.output_dir = root_dir / folder_name
        
        try:
            self.output_dir.mkdir(parents=True, exist_ok=True)
            logger.info(f"Created output directory: {self.output_dir}")
            return self.output_dir
        except Exception as e:
            logger.error(f"Failed to create output directory: {e}")
            # Fallback
            fallback = root_dir / f"Analysis_{int(now.timestamp())}"
            fallback.mkdir(exist_ok=True)
            return fallback

    def export_reports(self, clean_df: pd.DataFrame, stats_df: pd.DataFrame,
                      champion_df: pd.DataFrame, top_cells_df: pd.DataFrame,
                      yield_df: pd.DataFrame, batch_map_df: pd.DataFrame,
                      comparisons: Dict[str, Any], img_paths: Dict[str, Path],
                      user_initials: str, data_folder_name: str,
                      hysteresis_df: Optional[pd.DataFrame] = None) -> None:
        """
        Export Excel, Word, and PowerPoint reports.
        
        Args:
            clean_df: Cleaned dataframe
            stats_df: Statistics dataframe
            champion_df: Champion cells dataframe
            top_cells_df: Top 10 cells dataframe
            yield_df: Yield distribution dataframe
            batch_map_df: Batch mapping dataframe
            comparisons: Statistical comparison results
            img_paths: Dictionary of plot paths
            user_initials: User initials
            data_folder_name: Name of data folder
            hysteresis_df: Optional hysteresis dataframe
        """
        if not self.output_dir:
            logger.error("Output directory not set.")
            return

        logger.info("Exporting reports...")

        # Check report types
        report_types = getattr(self.config, 'report_types', ('excel', 'word', 'pptx'))

        # 1. Excel Export
        if 'excel' in report_types:
            try:
                excel_path = self.output_dir / self.config.excel_data_name
                with pd.ExcelWriter(excel_path, engine='openpyxl') as writer:
                    clean_df.to_excel(writer, sheet_name='Cleaned Data', index=False)
                    stats_df.to_excel(writer, sheet_name='Statistics', index=False)
                    champion_df.to_excel(writer, sheet_name='Champions', index=False)
                    top_cells_df.to_excel(writer, sheet_name='Top 10', index=False)
                    yield_df.to_excel(writer, sheet_name='Yield', index=False)
                    batch_map_df.to_excel(writer, sheet_name='Batch Map', index=False)
                    
                    if hysteresis_df is not None and not hysteresis_df.empty:
                        hysteresis_df.to_excel(writer, sheet_name='Hysteresis', index=False)
                
                logger.info(f"✓ Excel exported: {excel_path}")
            except Exception as e:
                logger.error(f"Excel export failed: {e}")

        # 2. Word Export
        if 'word' in report_types:
            try:
                self._export_word(clean_df, stats_df, champion_df, top_cells_df,
                                yield_df, batch_map_df, img_paths, hysteresis_df)
            except Exception as e:
                logger.error(f"Word export failed: {e}", exc_info=True)

        # 3. PowerPoint Export
        if 'pptx' in report_types:
            try:
                self._export_powerpoint(stats_df, champion_df, top_cells_df, 
                                      yield_df, batch_map_df, comparisons, 
                                      img_paths, user_initials, data_folder_name,
                                      clean_df, hysteresis_df)
            except Exception as e:
                logger.error(f"PowerPoint export failed: {e}", exc_info=True)

    def _export_word(self, clean_df: pd.DataFrame, stats_df: pd.DataFrame,
                    champion_df: pd.DataFrame, top_cells_df: pd.DataFrame,
                    yield_df: pd.DataFrame, batch_map_df: pd.DataFrame,
                    img_paths: Dict[str, Path],
                    hysteresis_df: Optional[pd.DataFrame] = None) -> None:
        """Generate Word report."""
        doc = Document()
        
        # Title
        title = doc.add_heading('IV Batch Analysis Report', level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Add timestamp
        doc.add_paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        doc.add_page_break()

        # Section 1: Summary
        doc.add_heading('1. Executive Summary', level=1)
        doc.add_paragraph(f"Total batches analyzed: {len(stats_df)}")
        doc.add_paragraph(f"Total cells: {len(clean_df)}")
        
        # Section 2: Tables
        doc.add_heading('2. Statistical Tables', level=1)
        
        self._add_word_table(doc, batch_map_df, ['Batch', 'Folder'], "2.1 Batch & Folder Reference")

        cols = ['Batch', 'Count', 'Eff_Mean', 'Eff_Max', 'Voc_Mean', 'FF_Mean']
        cols = [c for c in cols if c in stats_df.columns]
        self._add_word_table(doc, stats_df, cols, "2.2 Statistical Summary of PV Parameters")

        champ_cols = [c for c in ['Batch', 'CellName', 'Eff', 'Voc', 'Jsc', 'FF', 'Rs'] if c in champion_df.columns]
        self._add_word_table(doc, champion_df, champ_cols, "2.3 Champion Cell Parameters")
        self._add_word_table(doc, top_cells_df, champ_cols, "2.4 Top 10 Highest Efficiency Cells")

        yield_cols = [c for c in yield_df.columns if c != 'Batch']
        valid_yield_cols = ['Batch'] + [c for c in yield_cols if yield_df[c].sum() > 0]
        self._add_word_table(doc, yield_df, valid_yield_cols, "2.5 Efficiency Yield Distribution (%)")

        # Advanced Analysis Tables
        if self.config.enable_advanced_analysis:
            if 'Rs_fitted' in clean_df.columns:
                phy_cols = ['Batch', 'Rs_fitted', 'Rsh_fitted', 'n', 'I0', 'IL', 'fit_R2']
                phy_means = clean_df.groupby('Batch')[['Rs_fitted', 'Rsh_fitted', 'n', 'I0', 'IL', 'fit_R2']].mean().reset_index()
                self._add_word_table(doc, phy_means, phy_cols, "2.6 Physics Model Parameters (Mean)")

            if hysteresis_df is not None and not hysteresis_df.empty and 'HI_Eff' in hysteresis_df.columns:
                hys_cols = ['Batch', 'HI_Eff', 'HI_Voc', 'HI_Jsc', 'HI_FF']
                hys_means = hysteresis_df.groupby('Batch')[['HI_Eff', 'HI_Voc', 'HI_Jsc', 'HI_FF']].mean().reset_index()
                self._add_word_table(doc, hys_means, hys_cols, "2.7 Hysteresis Indices (Mean %)")

            if 'has_s_shape' in clean_df.columns:
                anom_counts = clean_df.groupby('Batch')[['has_s_shape', 'has_kink']].sum().reset_index()
                self._add_word_table(doc, anom_counts, ['Batch', 'has_s_shape', 'has_kink'], "2.8 Anomaly Detection Summary (Count)")

        doc.add_heading('3. Visual Analysis', level=1)
        plot_order = [
            ('jv_curve', 'J-V Curves of Champion Cells'),
            ('voc_jsc', 'Voc vs. Jsc Correlation Analysis'),
            ('box', 'Distribution of Electrical Parameters'),
            ('hist', 'Efficiency Histogram'),
            ('trend', 'Batch Trend Analysis (Max/Mean/Median)'),
            ('yield', 'Efficiency Yield Stack'),
            ('combo_drivers', 'Key Efficiency Drivers (Correlations)'),
            ('resistance', 'Parasitic Resistance Analysis (Rs & Rsh)'),
            ('model_fitting', 'Physics Model Fitting & Extraction'),
            ('hysteresis', 'Hysteresis Analysis'),
            ('anomalies', 'Anomaly Detection (S-Shapes)')
        ]
        for key, title in plot_order:
            if key in img_paths and img_paths[key]:
                doc.add_heading(title, level=2)
                doc.add_picture(str(img_paths[key]), width=Inches(6.5))
        
        doc.save(self.output_dir / self.config.report_docx_name)
        logger.info(f"✓ Word exported: {self.output_dir / self.config.report_docx_name}")

    def _export_powerpoint(self, stats_df: pd.DataFrame, champion_df: pd.DataFrame,
                          top_cells_df: pd.DataFrame, yield_df: pd.DataFrame,
                          batch_map_df: pd.DataFrame, comparisons: Dict[str, Any],
                          img_paths: Dict[str, Path], user_initials: str, data_folder_name: str = "",
                          clean_df: Optional[pd.DataFrame] = None,
                          hysteresis_df: Optional[pd.DataFrame] = None) -> None:
        """Export report to PowerPoint presentation."""
        prs = Presentation()
        prs.slide_width = PptInches(13.33)
        prs.slide_height = PptInches(7.5)

        # Title Slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        # Use data folder name if provided, otherwise fallback to generic title
        if data_folder_name:
            slide.shapes.title.text = f"IV Analysis Report - {data_folder_name}"
        else:
            slide.shapes.title.text = "IV Measurement Analysis"
        slide.placeholders[1].text = f"Operator: {user_initials}\nDate: {datetime.now().strftime('%Y-%m-%d')}"

        # Summary Slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Executive Summary"
        tf = slide.shapes.placeholders[1].text_frame
        slide.shapes.title.text_frame.paragraphs[0].font.size = PptPt(24)

        if comparisons.get('Results'):
            for batch, res in comparisons['Results'].items():
                p = tf.add_paragraph()
                p.text = f"{batch} vs {comparisons['Control']}:"
                p.font.bold = True
                p.font.size = PptPt(16)

                if 'Eff' in res:
                    p2 = tf.add_paragraph()
                    p2.level = 1
                    p2.text = f"Efficiency {res['Eff']['dir']} {abs(res['Eff']['diff']):.2f}% ({res['Eff']['sig']})"
                    p2.font.size = PptPt(16)

        cols = ['Batch', 'Count', 'Eff_Mean', 'Eff_Max', 'Voc_Mean', 'FF_Mean']
        cols = [c for c in cols if c in stats_df.columns]
        champ_cols = [c for c in ['Batch', 'CellName', 'Eff', 'Voc', 'Jsc', 'FF', 'Rs'] if c in champion_df.columns]
        yield_cols = [c for c in yield_df.columns if c != 'Batch']
        valid_yield_cols = ['Batch'] + [c for c in yield_cols if yield_df[c].sum() > 0]

        self._add_pptx_table(prs, batch_map_df, ['Batch', 'Folder'], "Batch & Folder Reference")
        self._add_pptx_table(prs, stats_df, cols, "Statistical Summary of PV Parameters")
        self._add_pptx_table(prs, champion_df, champ_cols, "Champion Cell Parameters")
        self._add_pptx_table(prs, top_cells_df, champ_cols, "Top 10 Highest Efficiency Cells")
        self._add_pptx_table(prs, yield_df, valid_yield_cols, "Efficiency Yield Distribution (%)")

        # --- Advanced Analysis Tables (PPT) ---
        if self.config.enable_advanced_analysis and clean_df is not None:
            if 'Rs_fitted' in clean_df.columns:
                phy_means = clean_df.groupby('Batch')[['Rs_fitted', 'Rsh_fitted', 'n', 'I0', 'IL', 'fit_R2']].mean().reset_index()
                self._add_pptx_table(prs, phy_means, ['Batch', 'Rs_fitted', 'Rsh_fitted', 'n', 'I0', 'IL'], "Physics Model Parameters (Mean)")

            if hysteresis_df is not None and not hysteresis_df.empty and 'HI_Eff' in hysteresis_df.columns:
                hys_means = hysteresis_df.groupby('Batch')[['HI_Eff', 'HI_Voc', 'HI_Jsc', 'HI_FF']].mean().reset_index()
                self._add_pptx_table(prs, hys_means, ['Batch', 'HI_Eff', 'HI_Voc', 'HI_Jsc', 'HI_FF'], "Hysteresis Indices (Mean %)")
            
            if 'has_s_shape' in clean_df.columns:
                anom_counts = clean_df.groupby('Batch')[['has_s_shape', 'has_kink']].sum().reset_index()
                self._add_pptx_table(prs, anom_counts, ['Batch', 'has_s_shape', 'has_kink'], "Anomaly Detection Summary (Count)")

        plot_order = [
            ('jv_curve', 'J-V Curves of Champion Cells'),
            ('voc_jsc', 'Voc vs. Jsc Correlation Analysis'),
            ('box', 'Distribution of Electrical Parameters'),
            ('hist', 'Efficiency Histogram'),
            ('trend', 'Batch Trend Analysis (Max/Mean/Median)'),
            ('yield', 'Efficiency Yield Stack'),
            ('combo_drivers', 'Key Efficiency Drivers (Correlations)'),
            ('resistance', 'Parasitic Resistance Analysis (Rs & Rsh)'),
            ('model_fitting', 'Physics Model Fitting & Extraction'),
            ('hysteresis', 'Hysteresis Analysis'),
            ('anomalies', 'Anomaly Detection (S-Shapes)')
        ]
        
        for key, title in plot_order:
            if key in img_paths and img_paths[key]:
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                slide.shapes.title.text = title
                slide.shapes.title.text_frame.paragraphs[0].font.size = PptPt(24)
                slide.shapes.add_picture(str(img_paths[key]), PptInches(2.5), PptInches(1.2),
                                        height=PptInches(5.8))

        prs.save(self.output_dir / self.config.report_pptx_name)
        logger.info(f"✓ PowerPoint exported: {self.output_dir / self.config.report_pptx_name}")

    # --- TABLE GENERATION METHODS ---

    @staticmethod
    def _add_word_table(doc: Document, df: pd.DataFrame, columns: list, title: str) -> None:
        """Add formatted table to Word document."""
        doc.add_heading(title, level=2)
        if df.empty or not columns:
            return

        table = doc.add_table(rows=1, cols=len(columns))
        table.style = 'Table Grid'
        hdr_cells = table.rows[0].cells

        for i, col in enumerate(columns):
            header_text = str(col).replace('_', ' (') + ')' if '_' in str(col) else str(col)
            hdr_cells[i].text = header_text
            p = hdr_cells[i].paragraphs[0]
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.runs[0]
            run.font.bold = True
            run.font.size = Pt(9)
            hdr_cells[i].vertical_alignment = WD_ALIGN_VERTICAL.CENTER

        for _, row in df.iterrows():
            row_cells = table.add_row().cells
            for i, col in enumerate(columns):
                val = row.get(col, "")
                if isinstance(val, (int, float)):
                    txt = f"{val:.0f}" if 'Count' in col else f"{val:.2f}"
                else:
                    txt = str(val)

                row_cells[i].text = txt
                p = row_cells[i].paragraphs[0]
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                p.runs[0].font.size = Pt(9)
                row_cells[i].vertical_alignment = WD_ALIGN_VERTICAL.CENTER
        doc.add_paragraph()

    @staticmethod
    def _add_pptx_table(prs: Presentation, df: pd.DataFrame, columns: list, title: str) -> None:
        """Add formatted table to PowerPoint presentation."""
        if df.empty or not columns:
            return
        max_rows = 10
        total_rows = len(df)
        num_slides = math.ceil(total_rows / max_rows)

        for i in range(num_slides):
            start = i * max_rows
            end = min((i + 1) * max_rows, total_rows)
            chunk = df.iloc[start:end]

            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"{title} ({i + 1}/{num_slides})" if num_slides > 1 else title
            slide.shapes.title.text_frame.paragraphs[0].font.size = PptPt(24)

            rows, cols = len(chunk) + 1, len(columns)
            left, top, width, height = PptInches(0.5), PptInches(1.5), PptInches(12.0), PptInches(0.5 * rows)

            table = slide.shapes.add_table(rows, cols, left, top, width, height).table

            for c_idx, col_name in enumerate(columns):
                cell = table.cell(0, c_idx)
                header_text = str(col_name).replace('_', ' (') + ')' if '_' in str(col_name) else str(col_name)
                cell.text = header_text
                p = cell.text_frame.paragraphs[0]
                p.font.bold = True
                p.font.size = PptPt(12)
                p.alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for r_idx in range(len(chunk)):
                row_data = chunk.iloc[r_idx]
                for c_idx, col_name in enumerate(columns):
                    cell = table.cell(r_idx + 1, c_idx)
                    val = row_data.get(col_name, "")
                    if isinstance(val, (int, float)):
                        txt = f"{val:.0f}" if 'Count' in col_name else f"{val:.2f}"
                    else:
                        txt = str(val)

                    cell.text = txt
                    p = cell.text_frame.paragraphs[0]
                    p.font.size = PptPt(11)
                    p.alignment = PP_ALIGN.CENTER
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
